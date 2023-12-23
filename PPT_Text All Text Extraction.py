import io
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image
import pytesseract
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def extract_text_from_powerpoint(file_path):
    def convert_tableto_html(table_data):
        html_table = "<table>"
        for table in table_data:
            for row in table:
                html_table += "<tr>"
                for cell in row:
                    html_table += f"<td>{cell}</td>"
                html_table += "</tr>"
        html_table += "</table>"
        return html_table
    #ppt file parsing
    prs = Presentation(file_path)
    ppt_content = ''
    #iterate slides
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            #read chart data
            if shape.has_chart:
                chart = shape.chart
                ppt_content+="Chart Details"
                for series_number, series in enumerate(chart.series, start=1):
                    if(chart.has_title):
                        if chart.chart_title.has_text_frame:
                            print("Chart Title: "+chart.chart_title.text_frame.text)                            
                            ppt_content +="Chart Title: "+chart.chart_title.text_frame.text
                    print(f" Series Name: {series.name}")
                    ppt_content +=series.name
                    print(f"  Values: {series.values}")
                    ppt_content +=series.values
                    if(chart.has_legend):
                        legend = chart.legend
                        if legend:
                            try:      
                                placeholderPart=slide.placeholders.part
                                #some inbuild tooltip need to avode in your case may be you get extra my case these are only
                                toolTipText_list = ["Click to edit", "Second level", 
                                                    "Third level", "Fourth level", 
                                                    "Fifth level", "Master title style"]         
                                for i,target_part in enumerate(placeholderPart.rels._rels):
                                    if(i>=0):
                                        try:                                            
                                            if (not placeholderPart.rels._rels[target_part].target_part.content_type.endswith(".drawingml.chart+xml")):
                                                continue
                                            root=ET.fromstring(placeholderPart.rels._rels[target_part].target_part.blob)
                                            for element in root.iter():
                                                if element.text:
                                                    isText=True
                                                    for pattern in toolTipText_list:
                                                        if element.text.startswith(pattern):
                                                            isText=False
                                                            break
                                                    if isText:
                                                        ppt_content += element.text.strip() + " "
                                        except Exception as e:
                                            pass
                            except Exception as e:
                                pass

                    print()
            print(shape.shape_type)
            #Read Nornam Text
            if hasattr(shape, "text"):
                ppt_content += shape.text + "\n"
            #Get images from slide    
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_stream = io.BytesIO(image_bytes)
                img = Image.open(image_stream)
                format = img.format.lower()
                if (format in ("jpeg" ,"png" ,"gif" ,"bmp" ,"tiff" ,"jpg")):
                    #using pytesseract read image text
                    img_text = pytesseract.image_to_string(img)
                    ppt_content += (img_text)
            #Get Tables from same slide and shape        
            table_data = []
            if shape.has_table:
                table = shape.table
                rows = []
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                table_data.append(rows)
                ppt_content  += convert_tableto_html(table_data)
            #get smartArt text from xml here first check condtion and get xml the to convert in nodes after the text by node
            if shape.has_chart==False and shape.has_table==False and  shape.has_text_frame==False and (shape.shape_type == None):
                try:      
                    placeholderPart=slide.placeholders.part
                    #some inbuild tooltip need to avode in your case may be you get extra my case these are only
                    toolTipText_list = ["Click to edit", "Second level", 
                                        "Third level", "Fourth level", 
                                        "Fifth level", "Master title style"]         
                    for i,target_part in enumerate(placeholderPart.rels._rels):
                        if(i>=0):
                            try:
                                print(placeholderPart.rels._rels[target_part].target_part.content_type)
                                if (placeholderPart.rels._rels[target_part].target_part.content_type.startswith('image')):
                                    continue
                                if (not #(placeholderPart.rels._rels[target_part].target_part.content_type.endswith('diagramDrawing+xml'))  
                                        #placeholderPart.rels._rels[target_part].target_part.content_type.endswith("diagramData+xml")
                                          placeholderPart.rels._rels[target_part].target_part.content_type.endswith("chart+xml")
                                          
                                         ):
                                    continue
                                root=ET.fromstring(placeholderPart.rels._rels[target_part].target_part.blob)
                                for element in root.iter():
                                    if element.text:
                                        isText=True
                                        for pattern in toolTipText_list:
                                            if element.text.startswith(pattern):
                                                isText=False
                                                break
                                        if isText:
                                            ppt_content += element.text.strip() + " "
                            except Exception as e:
                                pass
                except Exception as e:
                    pass
    return ppt_content
                
if __name__ == "__main__":
    file_path = "pptTest.pptx"
    result = extract_text_from_powerpoint(file_path)
    print (result)
