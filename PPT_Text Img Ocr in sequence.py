from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image
import pytesseract
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def extract_text_from_powerpoint(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_stream = BytesIO(image_bytes)
                img = Image.open(image_stream)
                img_text = pytesseract.image_to_string(img)
                text += (img_text)
                
    return text    
    

if __name__ == "__main__":
    file_path = "pptTest.pptx"
    result = extract_text_from_powerpoint(file_path)
    print (result)